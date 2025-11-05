# multi_file_tools.py - 多格式文件解析工具（支持xlsx/txt/pptx/image/audio/video/pdf）
import os
from pathlib import Path
from typing import Dict, Any
import json
import base64
import io

# -------------------------- 支持的文件类型定义 --------------------------
EXCEL_EXT = [".xls", ".xlsx"]
PDF_EXT = [".pdf"]
TXT_EXT = [".txt"]
PPTX_EXT = [".pptx"]
IMAGE_EXT = [".jpg", ".jpeg", ".png", ".bmp"]
AUDIO_EXT = [".mp3", ".wav"]
VIDEO_EXT = [".mp4", ".avi"]

# -------------------------- 核心解析函数 --------------------------
def prepare_file_for_llm(file_path: str) -> Dict[str, Any]:
    """
    根据文件类型读取并准备内容给多模态 LLM，返回格式：
    {
        "file": 文件路径,
        "type": 文件类型（excel/text/pptx/image/audio/video/pdf/unknown）,
        "llm_input": LLM 可解析的内容（文本/JSON/Base64）
    }
    """
    path = Path(file_path)
    ext = path.suffix.lower()
    file_name = path.name
    result = {"file": str(path), "type": None, "llm_input": None}

    # 1. Excel 文件：读取表格转为 JSON（取前10行避免过长）
    if ext in EXCEL_EXT:
        result["type"] = "excel"
        try:
            import pandas as pd
            df = pd.read_excel(file_path)
            # 转为列表字典，限制前10行
            result["llm_input"] = df.head(10).to_dict(orient="records")
        except Exception as e:
            result["llm_input"] = f"Excel read error: {str(e)}"

    # 2. TXT 文件：读取前10000字符（避免大文件溢出）
    elif ext in TXT_EXT:
        result["type"] = "text"
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read(10000)
                if len(text) >= 10000:
                    text += "\n...[Truncated: over 10000 characters]"
            result["llm_input"] = text
        except Exception as e:
            result["llm_input"] = f"TXT read error: {str(e)}"

    # 3. PPTX 文件：提取幻灯片文本+统计信息
    elif ext in PPTX_EXT:
        result["type"] = "pptx"
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            ppt_content = []
            image_count = 0

            # 遍历所有幻灯片提取文本和图片数量
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    # 提取文本框内容
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    # 统计图片数量（不提取二进制）
                    if shape.shape_type == 13:  # 13 = 图片类型
                        image_count += 1

                ppt_content.append({
                    "slide_index": slide_idx,
                    "text": "\n".join(slide_text) if slide_text else "No text"
                })

            # 构建 LLM 输入（包含基本信息+幻灯片内容）
            result["llm_input"] = {
                "pptx_info": {
                    "file_name": file_name,
                    "slide_count": len(prs.slides),
                    "image_count": image_count
                },
                "slide_content": ppt_content
            }
        except Exception as e:
            result["llm_input"] = f"PPTX read error: {str(e)}"

    # 4. 图片文件：提取尺寸+转为 Base64（供视觉分析）
    elif ext in IMAGE_EXT:
        result["type"] = "image"
        try:
            from PIL import Image
            with Image.open(file_path) as img:
                # 图片基本信息
                img_info = {
                    "format": img.format,
                    "size": f"{img.width}x{img.height}",
                    "mode": img.mode
                }
                # 压缩图片并转为 Base64（限制 800x600 像素）
                img.thumbnail((800, 600))
                img_byte_arr = io.BytesIO()
                save_format = img.format if img.format else "JPEG"
                img.save(img_byte_arr, format=save_format)
                img_base64 = base64.b64encode(img_byte_arr.getvalue()).decode("utf-8")

            result["llm_input"] = {
                "image_info": img_info,
                "image_base64": img_base64
            }
        except Exception as e:
            result["llm_input"] = f"Image read error: {str(e)}"

    # 5. 音频文件（MP3/WAV）：提取信息+转为 Base64
    elif ext in AUDIO_EXT:
        result["type"] = "audio"
        try:
            audio_info = {}
            # 提取音频基本信息（时长、比特率等）
            if ext == ".wav":
                import wave
                with wave.open(file_path, "rb") as wf:
                    audio_info["channels"] = wf.getnchannels()
                    audio_info["sample_rate"] = wf.getframerate()
                    audio_info["duration_second"] = round(wf.getnframes() / wf.getframerate(), 1)
            else:  # MP3
                import mutagen
                mp3 = mutagen.File(file_path)
                audio_info["duration_second"] = round(mp3.info.length, 1)
                audio_info["bitrate_kbps"] = mp3.info.bitrate // 1000

            # 读取音频二进制（限制 5MB，避免数据过大）
            max_size = 5 * 1024 * 1024  # 5MB
            with open(file_path, "rb") as f:
                audio_data = f.read(max_size)
                if len(audio_data) >= max_size:
                    audio_info["warning"] = "Truncated to 5MB"
                audio_base64 = base64.b64encode(audio_data).decode("utf-8")

            result["llm_input"] = {
                "audio_info": audio_info,
                "audio_base64": audio_base64
            }
        except Exception as e:
            result["llm_input"] = f"Audio read error: {str(e)}"

    # 6. 视频文件（MP4/AVI）：提取关键帧+转为 Base64
    elif ext in VIDEO_EXT:
        result["type"] = "video"
        try:
            import cv2
            from PIL import Image
            # 打开视频并提取关键帧
            cap = cv2.VideoCapture(file_path)
            if not cap.isOpened():
                raise Exception("Failed to open video")

            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            fps = cap.get(cv2.CAP_PROP_FPS)
            key_frames = []
            step = max(1, int(fps * 5))  # 每 5 秒取 1 帧

            # 提取前 5 帧（避免数据量过大）
            for i in range(0, min(frame_count, step * 5), step):
                cap.set(cv2.CAP_PROP_POS_FRAMES, i)
                ret, frame = cap.read()
                if ret:
                    # 转为 RGB 格式（OpenCV 默认 BGR）
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    img = Image.fromarray(frame_rgb)
                    # 压缩并转为 Base64
                    img.thumbnail((640, 480))
                    img_byte_arr = io.BytesIO()
                    img.save(img_byte_arr, format="JPEG")
                    img_base64 = base64.b64encode(img_byte_arr.getvalue()).decode("utf-8")

                    key_frames.append({
                        "time_second": round(i / fps, 1),
                        "frame_base64": img_base64
                    })

            cap.release()
            # 构建 LLM 输入
            result["llm_input"] = {
                "video_info": {
                    "file_name": file_name,
                    "duration_second": round(frame_count / fps, 1),
                    "fps": round(fps, 1),
                    "key_frame_count": len(key_frames)
                },
                "key_frames": key_frames
            }
        except Exception as e:
            result["llm_input"] = f"Video read error: {str(e)}"

    # 7. PDF 文件：提取文本（前2000字符）
    elif ext in PDF_EXT:
        result["type"] = "pdf"
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            # 提取所有页面文本并拼接
            text = "\n".join([page.extract_text() or "" for page in reader.pages])
            # 限制前2000字符
            if len(text) > 2000:
                text = text[:2000] + "\n...[Truncated: over 2000 characters]"
            result["llm_input"] = text
        except Exception as e:
            result["llm_input"] = f"PDF read error: {str(e)}"

    # 8. 未知格式
    else:
        result["type"] = "unknown"
        result["llm_input"] = f"Unsupported file type: {ext}"

    return result

# -------------------------- 测试函数 --------------------------
if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python multi_file_tools.py <file_path>")
        exit(1)
    file_path = sys.argv[1]
    info = prepare_file_for_llm(file_path)
    print(f"File Type: {info['type']}")
    print(f"LLM Input Preview: {json.dumps(info['llm_input'], ensure_ascii=False)[:500]}...")
