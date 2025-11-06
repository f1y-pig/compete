from pathlib import Path
from typing import Dict, Any

def handle_pptx(file_path: str) -> Dict[str, Any]:
    """
    处理 PPTX 文件，提取幻灯片文本和基本信息
    :param file_path: PPTX 文件路径
    :return: dict -> {file, type, llm_input, error}
    """
    path = Path(file_path)
    result = {"file": str(path), "type": "pptx", "llm_input": None, "error": None}

    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        ppt_content = []
        image_count = 0

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                if shape.shape_type == 13:  # 图片类型
                    image_count += 1
            ppt_content.append({
                "slide_index": slide_idx,
                "text": "\n".join(slide_text) if slide_text else "No text"
            })

        result["llm_input"] = {
            "pptx_info": {
                "file_name": path.name,
                "slide_count": len(prs.slides),
                "image_count": image_count
            },
            "slide_content": ppt_content
        }
    except Exception as e:
        result["error"] = f"PPTX read error: {str(e)}"

    return result
