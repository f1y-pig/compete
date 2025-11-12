"""
PPTX 专用工具：将幻灯片完整展开为 Markdown，或按页返回结构化数据。
自动适配 multi_file_tools_server.py 的路径逻辑。
"""
import json
import os
from collections import defaultdict
from pathlib import Path
from pptx import Presentation
from mcp.server.fastmcp import FastMCP
from pydantic import Field

mcp = FastMCP()

# 🔧 与 multi_file_tools_server.py 保持一致的路径定义
PROJECT_ROOT = r"/Users/dengken/Desktop/数据挖掘比赛/compete"
TEST_DIR = r"/Users/dengken/Desktop/数据挖掘比赛/compete/OxyGent-main/test"


def _resolve_file_path(file_path: str) -> Path:
    """
    自动在 TEST_DIR / PROJECT_ROOT 中寻找文件。
    """
    path = Path(file_path)
    if not path.is_absolute():
        test_path = Path(TEST_DIR) / path
        project_path = Path(PROJECT_ROOT) / path
        if test_path.exists():
            return test_path
        elif project_path.exists():
            return project_path
        else:
            return path  # 保留原始路径
    return path


def _slide_to_md(slide, index: int) -> str:
    """提取当前页的所有文本、表格、图片统计，转成 Markdown。"""
    lines = [f"## Slide {index}"]

    text_chunks = []
    table_chunks = []
    pictures = 0

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_chunks.append(shape.text.strip())
        if shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells))
            table_chunks.append("\n".join(rows))
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            pictures += 1

    if text_chunks:
        lines.append("\n".join(text_chunks))
    if table_chunks:
        lines.append("\n\n".join(f"表格:\n{tbl}" for tbl in table_chunks))
    if pictures:
        lines.append(f"*Images:* {pictures}")

    if len(lines) == 1:
        lines.append("_No textual content on this slide._")
    return "\n\n".join(lines)


@mcp.tool(description="将 PPTX 转为 Markdown，供 LLM 直接阅读")
def pptx_to_markdown(
    file_path: str = Field(..., description="PPTX 文件绝对路径或相对 test/ 的路径"),
    max_slides: int = Field(default=20, description="最多展开的页面数量，默认 20")
) -> str:
    """PPTX → Markdown"""
    path = _resolve_file_path(file_path)
    if not path.exists():
        return f"Error: file not found -> {path}"

    try:
        prs = Presentation(path)
    except Exception as exc:
        return f"Error: cannot open pptx -> {exc}"

    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        if idx > max_slides:
            slides.append(f"## Slide {idx}\n_Skipped due to max_slides limit._")
            break
        slides.append(_slide_to_md(slide, idx))

    meta = f"# PPTX Summary\n- File: {path.name}\n- Slides: {len(prs.slides)}\n"
    return meta + "\n\n".join(slides)


@mcp.tool(description="返回 PPTX 的结构化摘要（JSON）")
def pptx_to_json(
    file_path: str = Field(...),
    include_tables: bool = Field(default=True, description="是否包含表格内容"),
    include_images: bool = Field(default=True, description="是否统计图片数量")
) -> str:
    """PPTX → JSON"""
    path = _resolve_file_path(file_path)
    if not path.exists():
        return json.dumps({"error": f"file not found: {str(path)}"}, ensure_ascii=False)

    try:
        prs = Presentation(path)
    except Exception as exc:
        return json.dumps({"error": f"cannot open pptx: {exc}"}, ensure_ascii=False)

    summary = {
        "file": path.name,
        "absolute_path": str(path),
        "slide_count": len(prs.slides),
        "slides": []
    }

    for idx, slide in enumerate(prs.slides, 1):
        slide_info = defaultdict(list)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_info["texts"].append(shape.text.strip())
            if include_tables and shape.has_table:
                table_rows = []
                for row in shape.table.rows:
                    table_rows.append([cell.text.strip() for cell in row.cells])
                slide_info["tables"].append(table_rows)
            if include_images and shape.shape_type == 13:
                slide_info["images"].append("image_placeholder")

        summary["slides"].append({
            "index": idx,
            **slide_info
        })

    return json.dumps(summary, ensure_ascii=False, indent=2)


# ✅ 新增路径调试接口：方便和 multi_file_tools_server 共用
@mcp.tool(description="调试 PPTX 文件路径解析")
def debug_pptx_path(file_input: str = Field(..., description="输入文件名或路径")):
    path = _resolve_file_path(file_input)
    return {
        "input": file_input,
        "resolved_path": str(path),
        "exists": path.exists(),
        "test_dir": TEST_DIR,
        "project_root": PROJECT_ROOT
    }


if __name__ == "__main__":
    print("Starting PPTX QA Tool...")
    print(f"PROJECT_ROOT: {PROJECT_ROOT}")
    print(f"TEST_DIR: {TEST_DIR}")
    mcp.run()
