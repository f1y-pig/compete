# mcp_service/file_tools.py
"""
文件处理工具模块，用于 MCP 服务调用
支持 Excel/TXT/PPTX/PDF/图片/音频/视频文件
"""

from pathlib import Path
from typing import Dict, Any
from datetime import datetime
import base64
import io
import json
import subprocess

import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
from PIL import Image

# 音频处理
import wave
import mutagen

# 视频处理
import cv2

def call_video_tool(tool_name: str, payload: dict) -> dict:
    """
    简单的 MCP client 封装，复用你调用 video_tools 的方式。
    例如使用 subprocess.run(["python", "mcp_servers/video_tools.py", tool_name, ...])
    """
    result = subprocess.run(
        ["python", "mcp_servers/video_tools.py", tool_name, json.dumps(payload, ensure_ascii=False)],
        capture_output=True,
        text=True,
        check=False
    )
    try:
        return json.loads(result.stdout.strip() or "{}")
    except json.JSONDecodeError:
        return {"error": result.stdout}

def call_image_tool(tool_name: str, payload: dict) -> dict:
    result = subprocess.run(
        ["python", "mcp_servers/image_tools.py", tool_name, json.dumps(payload, ensure_ascii=False)],
        capture_output=True,
        text=True,
        check=False,
    )
    try:
        return json.loads(result.stdout.strip() or "{}")
    except json.JSONDecodeError:
        return {"status": "error", "message": result.stdout}

    
# -------------------- 文件处理函数 --------------------
def handle_excel(file_path: str) -> Dict[str, Any]:
    try:
        df = pd.read_excel(file_path)
        data = df.to_dict(orient="records")
        return {"file": file_path, "type": "excel", "content": data, "error": None}
    except Exception as e:
        return {"file": file_path, "type": "excel", "content": None, "error": str(e)}


def handle_txt(file_path: str) -> Dict[str, Any]:
    path = Path(file_path)
    result = {"file": str(path), "type": "text", "llm_input": None, "error": None}
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            text = f.read(10000)
            if len(text) >= 10000:
                text += "\n...[Truncated: over 10000 characters]"
        result["llm_input"] = text
    except Exception as e:
        result["error"] = f"TXT read error: {str(e)}"
    return result


def handle_pptx(file_path: str) -> Dict[str, Any]:
    path = Path(file_path)
    result = {"file": str(path), "type": "pptx", "llm_input": None, "error": None}
    try:
        prs = Presentation(file_path)
        ppt_content, image_count = [], 0
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                if shape.shape_type == 13:
                    image_count += 1
            ppt_content.append({"slide_index": idx, "text": "\n".join(slide_text) or "No text"})
        result["llm_input"] = {
            "pptx_info": {"file_name": path.name, "slide_count": len(prs.slides), "image_count": image_count},
            "slide_content": ppt_content
        }
    except Exception as e:
        result["error"] = f"PPTX read error: {str(e)}"
    return result


def handle_pdf(file_path: str) -> Dict[str, Any]:
    try:
        reader = PdfReader(file_path)
        text = "\n".join([page.extract_text() or "" for page in reader.pages])
        return {"file": file_path, "type": "pdf", "content": text, "error": None}
    except Exception as e:
        return {"file": file_path, "type": "pdf", "content": None, "error": str(e)}


def handle_image(file_path: str) -> dict:
    """
    统一处理图片：
    1. 调用 video_tools.ocr_image_file 做 OCR；
    2. 调用 image_tools.describe_images_with_qwen 获取场景/人物/商品；
    3. 把两段摘要合并写入 llm_input。
    """
    image_path = str(Path(file_path).resolve())
    ocr_summary = "（未识别到文字）"
    vision_summary = "（未生成视觉概述）"

    # 1. OCR
    try:
        ocr_resp = call_video_tool(
            "ocr_image_file",
            {
                "image_path": image_path,
                "lang": "chi_sim+eng",
            },
        )
        if ocr_resp.get("status") == "success" and ocr_resp.get("recognized_text"):
            ocr_summary = ocr_resp["recognized_text"].strip()
        else:
            error_msg = ocr_resp.get("error") or ocr_resp.get("message") or "未知原因"
            ocr_summary = f"OCR失败：{error_msg}"
    except Exception as exc:
        ocr_summary = f"OCR调用异常：{exc}"

    # 2. 视觉理解（多图时可传列表）
    try:
        vision_resp = call_image_tool(
            "describe_images_with_qwen",
            {
                "image_urls": json.dumps([image_path], ensure_ascii=False),
                "question": "概括图片中的场景、人物、商品及显著特征",
            },
        )
        if vision_resp.get("status") == "success" and vision_resp.get("summary"):
            vision_summary = vision_resp["summary"].strip()
        else:
            error_msg = vision_resp.get("message") or vision_resp.get("error") or "未知原因"
            vision_summary = f"视觉理解失败：{error_msg}"
    except Exception as exc:
        vision_summary = f"视觉工具调用异常：{exc}"

    llm_input = (
        "【OCR文字】\n"
        f"{ocr_summary}\n\n"
        "【图像理解】\n"
        f"{vision_summary}"
    )

    return {
        "type": "image",
        "llm_input": llm_input,
    }

def handle_audio(file_path: str) -> Dict[str, Any]:
    path = Path(file_path)
    result = {"file": str(path), "type": "audio", "llm_input": None, "error": None}
    try:
        audio_info = {}
        if path.suffix.lower() == ".wav":
            with wave.open(file_path, "rb") as wf:
                audio_info["channels"] = wf.getnchannels()
                audio_info["sample_rate"] = wf.getframerate()
                audio_info["duration_second"] = round(wf.getnframes() / wf.getframerate(), 1)
        elif path.suffix.lower() == ".mp3":
            mp3 = mutagen.File(file_path)
            audio_info["duration_second"] = round(mp3.info.length, 1)
            audio_info["bitrate_kbps"] = mp3.info.bitrate // 1000

        max_size = 5 * 1024 * 1024
        with open(file_path, "rb") as f:
            audio_data = f.read(max_size)
            if len(audio_data) >= max_size:
                audio_info["warning"] = "Truncated to 5MB"
            audio_base64 = base64.b64encode(audio_data).decode("utf-8")
        result["llm_input"] = {"audio_info": audio_info, "audio_base64": audio_base64}
    except Exception as e:
        result["error"] = f"Audio read error: {str(e)}"
    return result


def handle_video(file_path: str) -> dict:
    """
    通过 video_tools 的 batch_ocr_key_frames 获取视频文字摘要。
    返回格式保持 { "type": "video", "llm_input": ... }
    """
    try:
        # 直接调用 video_tools.py 的工具
        proc = subprocess.run(
            [
                "python",
                "mcp_servers/video_tools.py",
                "batch_ocr_key_frames",
                "--video_path",
                file_path,
                "--interval_seconds",
                "5",
            ],
            capture_output=True,
            text=True,
            check=True,
        )
        ocr_output = json.loads(proc.stdout)
    except Exception as e:
        return {
            "type": "video",
            "llm_input": f"Error summarizing video: {e}",
        }

    if "error" in ocr_output:
        return {"type": "video", "llm_input": ocr_output["error"]}

    summary = {
        "duration": ocr_output.get("duration_seconds"),
        "interval": ocr_output.get("interval_seconds"),
        "key_frames": ocr_output.get("key_frames", []),
    }
    return {
        "type": "video",
        "llm_input": json.dumps(summary, ensure_ascii=False),
    }



# -------------------- 统一分发函数 --------------------
def prepare_file_for_llm(file_path: str) -> Dict[str, Any]:
    """根据文件后缀调用对应处理函数"""
    ext = Path(file_path).suffix.lower()
    if ext in [".xls", ".xlsx"]:
        return handle_excel(file_path)
    elif ext == ".txt":
        return handle_txt(file_path)
    elif ext == ".pptx":
        return handle_pptx(file_path)
    elif ext in [".jpg", ".jpeg", ".png", ".bmp"]:
        return handle_image(file_path)
    elif ext in [".mp3", ".wav"]:
        return handle_audio(file_path)
    elif ext in [".mp4", ".avi"]:
        return handle_video(file_path)
    elif ext == ".pdf":
        return handle_pdf(file_path)
    else:
        return {"file": file_path, "type": "unknown", "llm_input": None, "error": f"Unsupported file type: {ext}"}
